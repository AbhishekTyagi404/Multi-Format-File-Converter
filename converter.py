from flask import Flask, request, render_template, send_file, flash, redirect
from pdf2docx import Converter
from pptx import Presentation
from fpdf import FPDF
from PIL import Image
import pypandoc
import os
import werkzeug
import docx
import time
import threading

# ✅ Set Pandoc Path
os.environ["PYPANDOC_PANDOC"] = "/home/kritvkjt/pandoc"

app = Flask(__name__)
app.secret_key = "supersecretkey"

UPLOAD_FOLDER = "uploads"
ALLOWED_EXTENSIONS = {"pdf", "docx", "txt", "pptx", "jpg", "jpeg", "png"}
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10 MB Limit

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER

# ✅ Function to Clean Old Files (Deletes files older than 30 minutes)
def cleanup_old_files():
    while True:
        current_time = time.time()
        for file_name in os.listdir(UPLOAD_FOLDER):
            file_path = os.path.join(UPLOAD_FOLDER, file_name)
            if os.path.isfile(file_path) and current_time - os.path.getmtime(file_path) > 1800:  # 30 minutes
                os.remove(file_path)
        time.sleep(300)  # Run cleanup every 5 minutes

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        file = request.files["file"]
        convert_type = request.form.get("convert_to")
        file_path = os.path.join(UPLOAD_FOLDER, werkzeug.utils.secure_filename(file.filename))
        file.save(file_path)

        output_path = None  # ✅ Initialize output_path

        # ✅ PDF to DOCX
        if convert_type == "docx" and file.filename.lower().endswith(".pdf"):
            output_path = file_path.replace(".pdf", ".docx")
            cv = Converter(file_path)
            cv.convert(output_path, start=0, end=None)
            cv.close()

        # ✅ DOCX to PDF (DOCX → HTML → PDF Fix for Unicode Support)
        elif convert_type == "pdf" and file.filename.lower().endswith(".docx"):
            html_path = file_path.replace(".docx", ".html")
            pdf_path = file_path.replace(".docx", ".pdf")

            # Convert DOCX to HTML
            pypandoc.convert_file(file_path, "html", outputfile=html_path)

            # Convert HTML to PDF using FPDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)

            with open(html_path, "r", encoding="utf-8") as html_file:
                for line in html_file:
                    pdf.multi_cell(0, 10, line)

            pdf.output(pdf_path)
            output_path = pdf_path

        # ✅ TXT to PDF
        elif convert_type == "pdf" and file.filename.lower().endswith(".txt"):
            output_path = file_path.replace(".txt", ".pdf")
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            with open(file_path, "r", encoding="utf-8") as txt_file:
                pdf.multi_cell(0, 10, txt_file.read())
            pdf.output(output_path, "F")

        # ✅ TXT to DOCX
        elif convert_type == "docx" and file.filename.lower().endswith(".txt"):
            output_path = file_path.replace(".txt", ".docx")
            doc = docx.Document()
            with open(file_path, "r", encoding="utf-8") as txt_file:
                doc.add_paragraph(txt_file.read())
            doc.save(output_path)

        # ✅ PPTX to PDF (Fixed Encoding for Hindi)
        elif convert_type == "pdf" and file.filename.lower().endswith(".pptx"):
            output_path = file_path.replace(".pptx", ".pdf")
            prs = Presentation(file_path)
            pdf = FPDF()
            pdf.set_auto_page_break(auto=True, margin=15)

            # ✅ Use Unicode Font (DejaVu)
            pdf.add_font("DejaVu", "", font_path, uni=True)
            pdf.set_font("DejaVu", "", 12)

            for slide in prs.slides:
                pdf.add_page()
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.encode("utf-8", "ignore").decode("utf-8")  # ✅ Fix Encoding
                        pdf.multi_cell(0, 10, text)

            pdf.output(output_path)

        # ✅ JPG/PNG to PDF
        elif convert_type == "pdf" and file.filename.lower().endswith((".jpg", ".jpeg", ".png")):
            output_path = file_path.rsplit(".", 1)[0] + ".pdf"
            img = Image.open(file_path)
            img.convert("RGB").save(output_path)

        # ✅ Handle unsupported file types
        if output_path is None or not os.path.exists(output_path):
            flash("Unsupported file format or conversion type.", "danger")
            return redirect(request.url)

        return send_file(output_path, as_attachment=True)

    return render_template("index.html")

if __name__ == "__main__":
    cleanup_thread = threading.Thread(target=cleanup_old_files, daemon=True)
    cleanup_thread.start()
    app.run()
